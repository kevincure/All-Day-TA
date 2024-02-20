# This code takes all documents (pdf, docx, doc, txt, tex, ppt, pptx) and chunks them, as well as saving a txt version of each document
# Run this before you run EmbedDocuments.py or app.py
# You need an OpenAI key saved in APIkey.txt
# Note that if your PDFs are not searchable, this won't work - use a third party tool to convert them to txt or doc first.  You
#   can look at the "-originaltext.csv" file created here and scan real quick to see if the text looks corrupted for any of your docs


import os
from PyPDF2 import PdfReader 
import nltk
import pandas as pd
import re
import openai
import shutil
from pptx import Presentation
# you need to pip install python-docx, not docx
import docx
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

# Set the desired chunk size and overlap size
# chunk_size is how many tokens we will take in each block of text
# overlap_size is how much overlap. So 200, 100 gives you chunks of between the 1st and 200th word, the 100th and 300th, the 200 and 400th...
# These chunk size and overlap are based on rules of thumb published in this literature, but not super optimized
chunk_size = 512
overlap_size = 100

# load user settings and api key
def read_settings(file_name):
    settings = {}
    with open(file_name, "r") as f:
        for line in f:
            key, value = line.strip().split("=")
            settings[key] = value
    return settings
settings = read_settings("settings.txt")
classname = settings["classname"]
classdescription = settings["classdescription"]
# get API_key
with open("APIkey.txt", "r") as f:
    openai.api_key = f.read().strip()
# Create subfolders if needed
filedirectory = "Documents"
output_folder = "Textchunks"
output_folder_full_text = "Full Text"
os.makedirs(output_folder, exist_ok=True)
os.makedirs(output_folder_full_text, exist_ok=True)

# counts the number of tokens in a chunk of text
def count_tokens(text):
    tokens = nltk.word_tokenize(text)
    return len(tokens)

# creates an overall summary of each document
def create_summary(row, classname, classdescription):
    # Check if the token count is more than 20,000; if so, just use first 20k
    if row['Token Count'] > 20000:
        tokens = nltk.word_tokenize(row['Text'])[:20000]
        text = ' '.join(tokens)
    else:
        text = row['Text']
    send_to_gpt=[]
    instructions = ("Consider this text which is a reading, transcript, slides or handout for " + classname + ", a " + classdescription + ". Give a ONE SENTENCE summary of what this specific text is about, such as 'This is a transcript from a class session describing Darwinian evolution in the context of the finch, including the history of science' or 'This appears to be a handout covering example computations of Cournot equilibria, Bertrand equilibria and other models of imperfect competition'.")
    try:
        send_to_gpt.append({"role": "system", "content": instructions})
        send_to_gpt.append({"role": "user", "content": text})
        response = openai.ChatCompletion.create(
            messages=send_to_gpt,
            temperature=0.1,
            max_tokens=100,
            model="gpt-4-turbo-preview"
        )
        tokens_sent = response["usage"]["prompt_tokens"]
        tokens_sent2 = response["usage"]["completion_tokens"]
        print(f"GPT-4-Turbo responded. You used {tokens_sent} prompt and {tokens_sent2} completion tokens.")
        current_summary = response["choices"][0]["message"]["content"]
    except Exception as e:
        current_summary = ""
        print("Error calling GPT to generate summary")
    return current_summary

# Add a chunk of text, with its title, to the dataframe
def append_to_dataframe(df, filename, text, title):
    # Create a new DataFrame row with the title and text
    new_row = pd.DataFrame({"Title": [title], "Text": [text]})
    # Append the new row to the existing DataFrame
    updated_df = pd.concat([df, new_row], ignore_index=True)
    return updated_df


# Loop through all pdf, txt, tex, ppt, pptx, in the "documents" folder
for filename in os.listdir(filedirectory):
    # Create an empty DataFrame to store the text and title of each document
    df = pd.DataFrame(columns=["Title", "Text"])
    print("Loading " + filename)
    filepath = os.path.join(filedirectory, filename)
    # Extract the title from the filename (remove the file extension)
    title = os.path.splitext(filename)[0]
    try:
        if filename.endswith(".pdf"):
            # Open the PDF file in read-binary mode
            reader = PdfReader(filepath)
            # Extract the text from each page of the PDF
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"

        elif filename.endswith(".ppt") or filename.endswith(".pptx"):
            ppt = Presentation(filepath)
            text = ''
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text

        elif filename.endswith(".doc") or filename.endswith(".docx"):
            doc = docx.Document(filepath)
            # Convert the file to UTF-8 and extract the text
            text = ''
            for paragraph in doc.paragraphs:
                text += paragraph.text

        elif filename.endswith(".txt"):
            # Open the text file and read its contents
            with open(filepath, "r", encoding="utf-8") as file:
                text = file.read()

        elif filename.endswith(".tex"):
            # Use regular expressions to extract regular text from the LaTeX file
            with open(filepath, "r", encoding="utf-8") as file:
                text = file.read()
            # remove everything before begin{document}
            match = re.search(r'\\begin\{document\}', text)
            if match:
                text = text[match.start():]
            # Replace special characters
            text = text.replace('\\$', '$')
            text = text.replace('\\\\', '\n')  # Replace \\ with newline for paragraph breaks
            # Remove comments
            text = re.sub(r'%.*?\n', '', text)
            def replace_math_expression(match):
                # Remove $ or $$ signs but keep the expression
                return match.group(1)
            # Modified regular expression to match both $...$ and $$...$$
            text= re.sub(r'\${1,2}(.*?)\${1,2}', replace_math_expression, text)
            # Remove common LaTeX commands
            commands = [
                r'\\textbf{.*?}', r'\\textit{.*?}', r'\\emph{.*?}', r'\\underline{.*?}',  # Formatting
                r'\\cite{.*?}', r'\\ref{.*?}',  # References
                r'\\label{.*?}',  # Labels
            ]
            for command in commands:
                text = re.sub(command, '', text)
    except Exception as e:
        print(f"Error processing file {filename}: {e}")
        continue  # Skip to the next file

    # Add the text and title to the DataFrame
    df = append_to_dataframe(df, filename, text, title)
    # save the full text from each file to "Full Text" - we'll use this to write questions
    output_full_text_file = os.path.join(output_folder_full_text, title + "-fulltext.txt")
    with open(output_full_text_file, 'w', encoding='utf-8') as file:
        file.write(text)

    # Create one-line sentences for each document, restricting to the first 20k words for very long documents
    df['Token Count'] = df['Text'].apply(count_tokens)
    # Initialize an empty list to store GPT responses
    summaries = []
    # Process each row and get LLM summary
    for index, row in df.iterrows():
        document_summary = create_summary(row, classname, classdescription)
        print(document_summary)
        summaries.append(document_summary)
    # Assign the list of responses to a new column in the DataFrame
    df['Summary'] = summaries

    # create chunks of the length and overlap specified above from each file
    chunks = []
    summary_chunks = []
    for i, row in df.iterrows():
        # Tokenize the text for the current row
        tokens = nltk.word_tokenize(row['Text'])
        # Loop through the tokens and create overlapping chunks
        for j in range(0, len(tokens), chunk_size - overlap_size):
            # Get the start and end indices of the current chunk
            start = j
            end = j + chunk_size
            # Create the current chunk by joining the tokens within the start and end indices
            chunk = ' '.join(tokens[start:end])
            # Add the article title to the beginning of the chunk
            chunk_with_title = "Source: " + row['Title'] + ". " + chunk
            chunk_with_title_and_summary = "Source: " + row['Title'] + ". " + "Summary of Document: " + row['Summary'] + " " + chunk
            # Append the current chunk to the list of chunks, along with the corresponding title
            chunks.append([row['Title'], chunk_with_title_and_summary, row['Summary'], chunk_with_title])
    # Convert the list of chunks to a dataframe
    df_chunks = pd.DataFrame(chunks, columns=['Title', 'Text', 'Summary', 'Raw Text'])
    # Truncate the filename if it's too long, e.g., limit to 250 characters
    max_filename_length = 250
    if len(filename) > max_filename_length:
        filename = filename[:max_filename_length]
    # Remove the file extension from the filename
    filename_without_extension = os.path.splitext(filename)[0]
    # Save the df_chunks to the output_folder subfolder with the new file name
    output_file = os.path.join(output_folder, filename_without_extension + "-originaltext.csv")
    df_chunks.to_csv(output_file, encoding='utf-8', escapechar='\\', index=False)
    print("Saving " + filename)

# move files you've chunked to an old directory
destination_directory = 'Already Chopped Documents'
for filename in os.listdir(filedirectory):
    source_path = os.path.join(filedirectory, filename)
    destination_path = os.path.join(destination_directory, filename)
    # Move the file to the destination directory
    shutil.move(source_path, destination_path)
print(f"Moved chopped documents to old directory")





